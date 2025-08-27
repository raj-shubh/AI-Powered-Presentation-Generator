from langchain.prompts import PromptTemplate
from pptx import Presentation
from pptx.util import Pt
from langchain_groq import ChatGroq
from langchain_tavily import TavilySearch
from dotenv import load_dotenv
import os
from langchain_core.output_parsers import JsonOutputParser
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

groq_api_key = os.getenv("GROQ_API_KEY")

llm = ChatGroq(
    model="llama3-8b-8192",   
    temperature=0.8,
    api_key=groq_api_key    
)

tavily_api_key = os.getenv('TAVILY_API_KEY')
search = TavilySearch(max_results=10,tavily_api_key=tavily_api_key)


topic = input("Enter your topic: ")
web_results = search.run(topic)

parser = JsonOutputParser()

template1 = PromptTemplate(
    template="""You are an expert researcher and presentation writer.  
Your task is to create a professional, presentation-ready summary** on the topic "{topic}".  
Use both your knowledge and the following web data for accurate, up-to-date information:  
{web_data}

Structure the summary with:
- A concise Introduction that sets context and defines the scope.
- 3–5 Key Themes or Sections, each with clear explanations, real-world examples, statistics, or case studies.
- Where possible, highlight trends, challenges, and opportunities.
- A Conclusion / Key Takeaways that is memorable and actionable.

Style Guidelines:
- Keep sentences short, impactful, and presentation-friendly.
- prefer clarity over complexity.
- Emphasize **facts, insights, and storytelling** rather than generic statements.
- Ensure smooth logical flow from start - middle - conclusion.

Return the response as a structured summary in plain text paragraphs, ready to be converted into slides.""",
    input_variables=['topic','web_data']
)


template2 = PromptTemplate(
    template= """
You are a professional presentation designer.  
Convert the following summary into a well-structured slide deck outline for a PowerPoint presentation:

{summary}

{format_instructions}

Formatting Rules:
- Output must be valid JSON only (no explanations or notes).
- JSON format:
[
  {{"title": "Slide 1 Title", "bullets": ["Point 1", "Point 2", "Point 3"]}},
  ...
]

Slide Structure:
- Slide 1 - Title slide (topic title + short subtitle if useful).  
- Slide 2 - Overview / Agenda.  
- Slides 3–6  Key insights, themes, or arguments (each with 3–5 concise bullet points).  
- Last Slide  Conclusion or Takeaways (clear, actionable, inspiring).  

Style Guidelines for Bullets:
- 1–2 sentences max per bullet.  
- Use clear, impactful wording.  
- Prefer facts, stats, and insights over filler text.  
- Ensure bullet points are self-contained and presentation-ready (no placeholders, no ellipses).

  Return ONLY valid JSON. Do not include any explanation, notes, or text before or after the JSON. 
  Do not include any notes, comments, or conditional logic like "or omit if not needed".
  Every bullet must be a complete string. Do not include ellipses or placeholders.""",
    input_variables=['summary'],
    partial_variables={'format_instructions': parser.get_format_instructions()}
)

chain = template1 | llm | template2 | llm | parser

slide_outline = chain.invoke({'topic':topic, 'web_data':web_results})

# print(slide_outline)


def create_ppt(slide_data, filename="LLM_Presentation2.pptx"):
    prs = Presentation()

    for i, slide_info in enumerate(slide_data):
        # 1st slide: title & subtitle)
        if i == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_info["title"]
            subtitle.text = "\n".join(slide_info["bullets"])

             # Style title
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) 

            # Style subtitle
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        else:
            slide_layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = slide_info["title"]
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 0) 
            text_frame = content.text_frame
            text_frame.clear()

            for bullet in slide_info["bullets"]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(22)
                p.font.name = 'Calibri'
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.level = 0

    prs.save(filename)
    print(f"Presentation saved as {filename}")
create_ppt(slide_outline, filename="LLM_Presentation2.pptx")